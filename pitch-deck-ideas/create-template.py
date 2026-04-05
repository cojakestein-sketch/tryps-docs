#!/usr/bin/env python3
"""
Generate tryps-template.pptx from pandoc's default reference.
Applies Tryps branding: colors, fonts, backgrounds.
Run: python3 create-template.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import re

# --- Brand constants ---
TRYPS_DARK    = RGBColor(0x0F, 0x19, 0x23)   # #0f1923 - primary dark
TRYPS_BLUE    = RGBColor(0x1A, 0x73, 0xE8)   # #1a73e8 - primary blue
TRYPS_ACCENT  = RGBColor(0xE9, 0x45, 0x60)   # #e94560 - accent red
TRYPS_GRAY_BG = RGBColor(0xF8, 0xF9, 0xFA)   # #f8f9fa - light gray bg
TRYPS_TEXT     = RGBColor(0x1A, 0x1A, 0x1A)   # #1a1a1a - body text
TRYPS_MUTED    = RGBColor(0x6B, 0x72, 0x80)   # #6b7280 - muted text
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)

# Font family — Inter if installed, falls back to system sans
FONT_HEADING = "Inter"
FONT_BODY    = "Inter"


def set_solid_bg(slide_or_layout, color):
    """Set a solid background color on a slide layout."""
    bg = slide_or_layout.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def style_placeholder(ph, font_name=None, font_size=None, font_color=None,
                      bold=None, alignment=None):
    """Style all paragraphs in a placeholder's default text frame."""
    if not ph.has_text_frame:
        return
    for para in ph.text_frame.paragraphs:
        if alignment is not None:
            para.alignment = alignment
        for run in para.runs:
            if font_name:
                run.font.name = font_name
            if font_size:
                run.font.size = font_size
            if font_color:
                run.font.color.rgb = font_color
            if bold is not None:
                run.font.bold = bold

    # Also set the default paragraph font via XML for new content
    for para in ph.text_frame.paragraphs:
        defRPr = para._p.find(qn('a:pPr'))
        if defRPr is None:
            defRPr = para._p.makeelement(qn('a:pPr'), {})
            para._p.insert(0, defRPr)

        rPr = defRPr.find(qn('a:defRPr'))
        if rPr is None:
            rPr = defRPr.makeelement(qn('a:defRPr'), {})
            defRPr.append(rPr)

        if font_name:
            # Set latin font
            latin = rPr.find(qn('a:latin'))
            if latin is None:
                latin = rPr.makeelement(qn('a:latin'), {})
                rPr.append(latin)
            latin.set('typeface', font_name)

        if font_size:
            rPr.set('sz', str(int(font_size.pt * 100)))

        if font_color:
            solidFill = rPr.find(qn('a:solidFill'))
            if solidFill is None:
                solidFill = rPr.makeelement(qn('a:solidFill'), {})
                rPr.insert(0, solidFill)
            else:
                solidFill.clear()
            srgb = solidFill.makeelement(qn('a:srgbClr'), {
                'val': str(font_color)
            })
            solidFill.append(srgb)

        if bold is not None:
            rPr.set('b', '1' if bold else '0')


def update_theme_colors(prs):
    """Update the theme XML to use Tryps brand colors."""
    for rel in prs.slide_masters[0].part.rels.values():
        if 'theme' in rel.reltype:
            theme_part = rel.target_part
            theme_xml = etree.fromstring(theme_part.blob)

            # Find clrScheme
            clr_scheme = theme_xml.find('.//' + qn('a:clrScheme'))
            if clr_scheme is not None:
                clr_scheme.set('name', 'Tryps')
                color_map = {
                    'dk1':     '0F1923',  # Dark 1 → Tryps dark
                    'lt1':     'FFFFFF',  # Light 1 → White
                    'dk2':     '1A1A1A',  # Dark 2 → Text
                    'lt2':     'F8F9FA',  # Light 2 → Light gray
                    'accent1': '1A73E8',  # Accent 1 → Tryps blue
                    'accent2': 'E94560',  # Accent 2 → Tryps accent red
                    'accent3': '0F1923',  # Accent 3 → Dark
                    'accent4': '6B7280',  # Accent 4 → Muted
                    'accent5': '34D399',  # Accent 5 → Green (for positive)
                    'accent6': 'F59E0B',  # Accent 6 → Amber (for warning)
                    'hlink':   '1A73E8',  # Hyperlink → Blue
                    'folHlink':'6B7280',  # Followed link → Muted
                }
                for elem_name, hex_color in color_map.items():
                    elem = clr_scheme.find(qn(f'a:{elem_name}'))
                    if elem is not None:
                        srgb = elem.find(qn('a:srgbClr'))
                        if srgb is not None:
                            srgb.set('val', hex_color)
                        else:
                            # Replace sysClr or whatever with srgbClr
                            for child in list(elem):
                                elem.remove(child)
                            new_srgb = etree.SubElement(elem, qn('a:srgbClr'))
                            new_srgb.set('val', hex_color)

            # Update font scheme
            font_scheme = theme_xml.find('.//' + qn('a:fontScheme'))
            if font_scheme is not None:
                font_scheme.set('name', 'Tryps')
                for font_type in ['majorFont', 'minorFont']:
                    font_elem = font_scheme.find(qn(f'a:{font_type}'))
                    if font_elem is not None:
                        latin = font_elem.find(qn('a:latin'))
                        if latin is not None:
                            latin.set('typeface', FONT_HEADING if font_type == 'majorFont' else FONT_BODY)

            # Write back the modified theme XML
            theme_part._blob = etree.tostring(theme_xml, xml_declaration=True,
                                               encoding='UTF-8', standalone=True)
            break


def main():
    prs = Presentation('/tmp/pandoc-default-reference.pptx')

    # Set 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Update theme colors and fonts
    update_theme_colors(prs)

    layouts = {layout.name: layout for layout in prs.slide_layouts}

    # --- Layout 0: Title Slide (dark background) ---
    title_layout = layouts['Title Slide']
    set_solid_bg(title_layout, TRYPS_DARK)

    for ph in title_layout.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:  # Title
            style_placeholder(ph, FONT_HEADING, Pt(44), WHITE, bold=True,
                              alignment=PP_ALIGN.CENTER)
        elif idx == 1:  # Subtitle
            style_placeholder(ph, FONT_BODY, Pt(22), RGBColor(0xAA, 0xAA, 0xBB),
                              alignment=PP_ALIGN.CENTER)

    # --- Layout 1: Title and Content (white background) ---
    content_layout = layouts['Title and Content']
    set_solid_bg(content_layout, WHITE)

    for ph in content_layout.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:  # Title
            style_placeholder(ph, FONT_HEADING, Pt(32), TRYPS_DARK, bold=True)
        elif idx == 1:  # Content
            style_placeholder(ph, FONT_BODY, Pt(20), TRYPS_TEXT)

    # --- Layout 2: Section Header (dark background) ---
    section_layout = layouts['Section Header']
    set_solid_bg(section_layout, TRYPS_DARK)

    for ph in section_layout.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:  # Title
            style_placeholder(ph, FONT_HEADING, Pt(40), WHITE, bold=True)
        elif idx == 1:  # Subtitle
            style_placeholder(ph, FONT_BODY, Pt(20), RGBColor(0xAA, 0xAA, 0xBB))

    # --- Layout 3: Two Content (white background) ---
    two_content_layout = layouts['Two Content']
    set_solid_bg(two_content_layout, WHITE)

    for ph in two_content_layout.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            style_placeholder(ph, FONT_HEADING, Pt(32), TRYPS_DARK, bold=True)
        elif idx in (1, 2):
            style_placeholder(ph, FONT_BODY, Pt(18), TRYPS_TEXT)

    # --- Layout 4: Comparison (white background) ---
    comparison_layout = layouts['Comparison']
    set_solid_bg(comparison_layout, WHITE)

    for ph in comparison_layout.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            style_placeholder(ph, FONT_HEADING, Pt(32), TRYPS_DARK, bold=True)
        else:
            style_placeholder(ph, FONT_BODY, Pt(18), TRYPS_TEXT)

    # --- Layout 6: Blank (white) ---
    blank_layout = layouts['Blank']
    set_solid_bg(blank_layout, WHITE)

    # --- Layout 7: Content with Caption (light gray background) ---
    caption_layout = layouts['Content with Caption']
    set_solid_bg(caption_layout, TRYPS_GRAY_BG)

    for ph in caption_layout.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            style_placeholder(ph, FONT_HEADING, Pt(28), TRYPS_DARK, bold=True)
        elif idx == 1:
            style_placeholder(ph, FONT_BODY, Pt(18), TRYPS_TEXT)
        elif idx == 2:
            style_placeholder(ph, FONT_BODY, Pt(16), TRYPS_MUTED)

    # --- Remove all slides (template should have no content slides) ---
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Save
    out_path = '/Users/jakestein/tryps-docs/pitch-deck-ideas/tryps-template.pptx'
    prs.save(out_path)
    print(f"Created {out_path}")
    print()
    print("Layouts branded:")
    print("  Title Slide      → dark navy bg, white Inter 44pt")
    print("  Section Header   → dark navy bg, white Inter 40pt")
    print("  Title and Content → white bg, dark Inter 32pt / 20pt")
    print("  Two Content      → white bg, two columns")
    print("  Comparison       → white bg, comparison columns")
    print("  Content w/ Caption → light gray bg")
    print("  Blank            → white bg")
    print()
    print("Theme colors: Tryps dark (#0f1923), blue (#1a73e8), accent (#e94560)")
    print("Fonts: Inter (heading + body)")


if __name__ == '__main__':
    main()
